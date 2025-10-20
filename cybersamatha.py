import os
import glob
import argparse
from typing import List, Dict, Any
import hashlib
from datetime import datetime

# Import required packages
try:
    import google.generativeai as genai
    import chromadb
    from chromadb.utils import embedding_functions
    from sentence_transformers import SentenceTransformer
    import PyPDF2
    from docx import Document
    import pptx
    import json
    import yaml
    from dotenv import load_dotenv
except ImportError as e:
    print(f"❌ Missing required package: {e}")
    print("📦 Please install required packages:")
    print("pip install google-generativeai chromadb sentence-transformers pypdf2 python-docx python-pptx pyyaml python-dotenv")
    exit(1)

class CyberSamathaRAG:
    def __init__(self, data_path: str = "data", chroma_path: str = "chroma_db", 
                 embedding_model: str = "all-MiniLM-L6-v2"):
        self.data_path = data_path
        self.chroma_path = chroma_path
        self.chroma_client = None
        self.collection = None
        self.embedding_model_name = embedding_model
        self.sentence_transformer = None
        self.genai_model = None
        
        # Create directories if they don't exist
        os.makedirs(self.data_path, exist_ok=True)
        os.makedirs(self.chroma_path, exist_ok=True)
        
        # Initialize components
        self._setup_local_embeddings()
        self._setup_gemini()
        self._setup_vector_db()
        
    def _setup_local_embeddings(self):
        """Initialize local sentence transformer model for embeddings"""
        try:
            print(f"📥 Loading local embedding model: {self.embedding_model_name}")
            self.sentence_transformer = SentenceTransformer(self.embedding_model_name)
            print("✅ Local embedding model loaded successfully")
        except Exception as e:
            print(f"❌ Failed to load embedding model: {e}")
            print("💡 Downloading model for first-time use...")
            try:
                self.sentence_transformer = SentenceTransformer(self.embedding_model_name)
                print("✅ Model downloaded and loaded successfully")
            except Exception as e2:
                raise ValueError(f"❌ Cannot initialize embedding model: {e2}")
        
    def _setup_gemini(self):
        """Initialize Gemini API for answer generation only"""
        load_dotenv()
        
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            print("⚠️  GEMINI_API_KEY not found - running in embedding-only mode")
            print("💡 You can still index documents and search, but won't get AI-generated answers")
            print("   To enable AI answers, set GEMINI_API_KEY in .env file")
            print("   Get your key from: https://aistudio.google.com/app/apikey")
            self.genai_model = None
            return
        
        try:
            genai.configure(api_key=api_key)
            self.genai_model = genai.GenerativeModel('gemini-pro')
            print("✅ Gemini API configured for answer generation")
        except Exception as e:
            print(f"⚠️  Failed to configure Gemini API: {e}")
            print("   Running in embedding-only mode")
            self.genai_model = None
        
    def _setup_vector_db(self):
        """Initialize ChromaDB with local embedding function"""
        try:
            # Use PersistentClient for disk storage
            self.chroma_client = chromadb.PersistentClient(path=self.chroma_path)
            
            # Create custom embedding function using sentence-transformers
            class LocalEmbeddingFunction(embedding_functions.EmbeddingFunction):
                def __init__(self, model):
                    self.model = model
                
                def __call__(self, input: List[str]) -> List[List[float]]:
                    embeddings = self.model.encode(input, convert_to_numpy=True)
                    return embeddings.tolist()
            
            self.embedding_function = LocalEmbeddingFunction(self.sentence_transformer)
            
            # Create or get collection
            try:
                self.collection = self.chroma_client.get_collection(
                    name="cybersamatha_docs",
                    embedding_function=self.embedding_function
                )
                print("✅ Loaded existing vector database")
            except Exception:
                self.collection = self.chroma_client.create_collection(
                    name="cybersamatha_docs",
                    embedding_function=self.embedding_function,
                    metadata={"hnsw:space": "cosine"}
                )
                print("✅ Created new vector database")
                
        except Exception as e:
            print(f"❌ Error setting up vector database: {e}")
            raise
    
    def _read_file(self, file_path: str) -> str:
        """Read various file formats and extract text"""
        try:
            if file_path.endswith('.pdf'):
                return self._read_pdf(file_path)
            elif file_path.endswith('.docx'):
                return self._read_docx(file_path)
            elif file_path.endswith('.pptx'):
                return self._read_pptx(file_path)
            elif file_path.endswith('.json'):
                return self._read_json(file_path)
            elif file_path.endswith('.yaml') or file_path.endswith('.yml'):
                return self._read_yaml(file_path)
            else:  # .txt, .md, etc.
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except Exception as e:
            print(f"❌ Error reading {file_path}: {str(e)}")
            return ""
    
    def _read_pdf(self, file_path: str) -> str:
        """Extract text from PDF files"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
        except Exception as e:
            print(f"❌ Error reading PDF {file_path}: {e}")
        return text
    
    def _read_docx(self, file_path: str) -> str:
        """Extract text from Word documents"""
        try:
            doc = Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            print(f"❌ Error reading DOCX {file_path}: {e}")
            return ""
    
    def _read_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        text = ""
        try:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"❌ Error reading PPTX {file_path}: {e}")
        return text
    
    def _read_json(self, file_path: str) -> str:
        """Extract text from JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        except Exception as e:
            print(f"❌ Error reading JSON {file_path}: {e}")
            return ""
    
    def _read_yaml(self, file_path: str) -> str:
        """Extract text from YAML files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
                return yaml.dump(data, default_flow_style=False)
        except Exception as e:
            print(f"❌ Error reading YAML {file_path}: {e}")
            return ""
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if not text or len(text.strip()) == 0:
            return []
            
        if len(text) <= chunk_size:
            return [text.strip()]
        
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            if end < len(text):
                # Try to break at sentence end
                last_period = text.rfind('. ', start, end)
                last_newline = text.rfind('\n', start, end)
                
                if last_period != -1 and last_period > start + chunk_size // 2:
                    end = last_period + 1
                elif last_newline != -1 and last_newline > start + chunk_size // 2:
                    end = last_newline + 1
                    
            chunk = text[start:end].strip()
            if chunk:  # Only add non-empty chunks
                chunks.append(chunk)
            start = end - chunk_overlap
            
            # Prevent infinite loop
            if start >= len(text):
                break
        
        return chunks
    
    def _get_file_hash(self, file_path: str) -> str:
        """Generate hash of file for change detection"""
        try:
            hasher = hashlib.md5()
            with open(file_path, 'rb') as f:
                buf = f.read()
                hasher.update(buf)
            return hasher.hexdigest()
        except Exception as e:
            print(f"❌ Error hashing file {file_path}: {e}")
            return ""
    
    def index_documents(self, force_reindex: bool = False):
        """Index all documents in the data folder"""
        print("🔍 Scanning for documents...")
        
        # Supported file extensions
        extensions = ['*.txt', '*.md', '*.json', '*.yaml', '*.yml', 
                     '*.pdf', '*.docx', '*.pptx']
        
        all_files = []
        for ext in extensions:
            all_files.extend(glob.glob(os.path.join(self.data_path, '**', ext), recursive=True))
            all_files.extend(glob.glob(os.path.join(self.data_path, '**', ext.upper()), recursive=True))
        
        if not all_files:
            print("📁 No documents found in the 'data' folder")
            print("💡 Add some documents (PDF, DOCX, TXT, etc.) to the 'data' folder and run again")
            return
        
        print(f"📁 Found {len(all_files)} documents")
        
        processed_count = 0
        total_chunks = 0
        skipped_count = 0
        
        for file_path in all_files:
            try:
                file_hash = self._get_file_hash(file_path)
                file_name = os.path.basename(file_path)
                relative_path = os.path.relpath(file_path, self.data_path)
                
                if not file_hash:
                    continue
                
                # Check if file already indexed
                if not force_reindex:
                    try:
                        existing = self.collection.get(
                            where={"file_path": relative_path},
                            include=["metadatas"]
                        )
                        if existing["metadatas"] and existing["metadatas"][0].get("file_hash") == file_hash:
                            print(f"⏭️  Skipping unchanged: {relative_path}")
                            skipped_count += 1
                            continue
                        elif existing["ids"]:
                            # Delete old chunks before reindexing
                            print(f"🔄 Updating: {relative_path}")
                            self.collection.delete(ids=existing["ids"])
                    except Exception:
                        pass  # Collection might be empty
                
                # Read and process file
                print(f"📖 Processing: {relative_path}")
                content = self._read_file(file_path)
                
                if not content or not content.strip():
                    print(f"⚠️  Empty or unreadable content: {relative_path}")
                    continue
                
                # Split into chunks
                chunks = self._chunk_text(content)
                
                if not chunks:
                    print(f"⚠️  No chunks created: {relative_path}")
                    continue
                
                # Prepare documents for indexing
                documents = []
                metadatas = []
                ids = []
                
                for i, chunk in enumerate(chunks):
                    chunk_id = f"{relative_path}_chunk_{i}_{hashlib.md5(chunk.encode()).hexdigest()[:8]}"
                    documents.append(chunk)
                    metadatas.append({
                        "file_path": relative_path,
                        "file_name": file_name,
                        "chunk_index": i,
                        "file_hash": file_hash,
                        "timestamp": datetime.now().isoformat(),
                        "chunk_size": len(chunk),
                        "embedding_model": self.embedding_model_name
                    })
                    ids.append(chunk_id)
                
                # Add to vector database (embeddings computed automatically by Chroma)
                self.collection.add(
                    documents=documents,
                    metadatas=metadatas,
                    ids=ids
                )
                
                processed_count += 1
                total_chunks += len(chunks)
                print(f"✅ Indexed {len(chunks)} chunks from {relative_path}")
                
            except Exception as e:
                print(f"❌ Error processing {file_path}: {str(e)}")
        
        print("\n" + "="*60)
        if processed_count > 0:
            print(f"🎉 Indexing complete!")
            print(f"   📝 Processed: {processed_count} files")
            print(f"   📦 Created: {total_chunks} chunks")
            if skipped_count > 0:
                print(f"   ⏭️  Skipped: {skipped_count} unchanged files")
            print(f"   🧠 Using: {self.embedding_model_name} (local)")
        else:
            print("ℹ️  No new files to index")
            if skipped_count > 0:
                print(f"   ⏭️  {skipped_count} files already indexed")
        print("="*60)
    
    def search_documents(self, query: str, n_results: int = 5) -> List[Dict]:
        """Search for relevant documents using local embeddings"""
        try:
            results = self.collection.query(
                query_texts=[query],
                n_results=n_results,
                include=["documents", "metadatas", "distances"]
            )
            
            search_results = []
            if results["documents"] and results["documents"][0]:
                for doc, meta, distance in zip(
                    results["documents"][0],
                    results["metadatas"][0],
                    results["distances"][0]
                ):
                    search_results.append({
                        "content": doc,
                        "metadata": meta,
                        "score": 1 - distance,  # Convert distance to similarity score
                        "file_name": meta.get("file_name", "Unknown"),
                        "file_path": meta.get("file_path", "Unknown")
                    })
            
            return search_results
        except Exception as e:
            print(f"❌ Search error: {e}")
            return []
    
    def ask_question(self, question: str, context_chunks: int = 5, show_sources: bool = True) -> str:
        """Ask a question using RAG with local embeddings + Gemini reasoning"""
        print("🔍 Searching knowledge base...")
        
        # Search for relevant context using local embeddings
        search_results = self.search_documents(question, n_results=context_chunks)
        
        if not search_results:
            return "❌ No relevant information found in the knowledge base.\n\n💡 Try indexing some documents first using: python cybersamatha.py --index"
        
        # Build context from search results
        context_parts = []
        sources = []
        for i, result in enumerate(search_results, 1):
            source_file = result["metadata"]["file_path"]
            relevance = result['score']
            context_parts.append(f"[Document {i}: {source_file} | Relevance: {relevance:.2f}]\n{result['content']}")
            sources.append(f"  {i}. {source_file} (relevance: {relevance:.2%})")
        
        context = "\n\n".join(context_parts)
        
        # If Gemini is not available, return search results only
        if self.genai_model is None:
            answer = "🔍 **Search Results** (AI answer generation disabled - no GEMINI_API_KEY)\n\n"
            answer += f"Found {len(search_results)} relevant documents:\n\n"
            for i, result in enumerate(search_results, 1):
                answer += f"**Source {i}: {result['file_name']}** (relevance: {result['score']:.2%})\n"
                answer += f"{result['content'][:500]}...\n\n"
            return answer
        
        # Create prompt with context for Gemini
        prompt = f"""Based on the following cybersecurity documentation, please answer the question thoroughly and accurately.

CONTEXT DOCUMENTS:
{context}

QUESTION: {question}

INSTRUCTIONS:
- Provide a comprehensive answer based strictly on the context provided
- If the context doesn't contain enough information, acknowledge this limitation
- Include specific references to source documents when possible
- Format your response in a clear, organized manner
- Focus on actionable cybersecurity insights
- Be concise but thorough"""

        try:
            print("🤖 Generating AI answer...")
            response = self.genai_model.generate_content(prompt)
            answer = response.text
            
            if show_sources:
                answer += "\n\n📚 **Sources:**\n" + "\n".join(sources)
            
            return answer
        except Exception as e:
            error_msg = f"❌ Error generating AI response: {str(e)}\n\n"
            error_msg += "📋 **Retrieved Context:**\n" + "\n".join(sources)
            return error_msg
    
    def get_collection_stats(self) -> Dict[str, Any]:
        """Get statistics about the vector database"""
        try:
            count = self.collection.count()
            return {
                "total_chunks": count,
                "collection_name": "cybersamatha_docs",
                "database_path": self.chroma_path,
                "embedding_model": self.embedding_model_name,
                "gemini_enabled": self.genai_model is not None
            }
        except Exception as e:
            return {"error": str(e)}
    
    def interactive_chat(self):
        """Start interactive chat session"""
        stats = self.get_collection_stats()
        total_chunks = stats.get("total_chunks", 0)
        
        print("\n" + "="*70)
        print("🤖 CyberSamatha RAG System - Ready!")
        print("="*70)
        print(f"📊 Knowledge Base: {total_chunks} document chunks")
        print(f"🧠 Embedding Model: {self.embedding_model_name} (local)")
        print(f"🤖 AI Answers: {'Enabled (Gemini)' if stats.get('gemini_enabled') else 'Disabled'}")
        print("💬 Ask questions about your cybersecurity documents")
        print("⏹️  Commands: 'quit'/'exit' (end), 'stats' (show stats)")
        print("="*70 + "\n")
        
        while True:
            try:
                question = input("💬 Question: ").strip()
                
                if question.lower() in ['quit', 'exit', 'bye', 'q']:
                    print("👋 Goodbye!")
                    break
                
                if not question:
                    continue
                
                if question.lower() in ['stats', 'status', 'info']:
                    stats = self.get_collection_stats()
                    print(f"\n📊 **Database Statistics:**")
                    print(f"   • Total chunks: {stats.get('total_chunks', 0)}")
                    print(f"   • Embedding model: {stats.get('embedding_model', 'N/A')}")
                    print(f"   • AI answers: {'✅ Enabled' if stats.get('gemini_enabled') else '❌ Disabled'}")
                    print(f"   • Database path: {stats.get('database_path', 'N/A')}\n")
                    continue
                
                answer = self.ask_question(question)
                print(f"\n🤖 **Answer:**\n{answer}\n")
                print("-" * 80)
                
            except KeyboardInterrupt:
                print("\n👋 Goodbye!")
                break
            except Exception as e:
                print(f"❌ Error: {str(e)}")

def main():
    parser = argparse.ArgumentParser(
        description="🤖 CyberSamatha RAG System - Local Embeddings + Gemini Reasoning",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python cybersamatha.py --index                    # Index documents (local embeddings)
  python cybersamatha.py --index --force            # Force reindex all documents
  python cybersamatha.py --question "What is XSS?"  # Ask a question
  python cybersamatha.py --stats                    # Show database statistics
  python cybersamatha.py                            # Start interactive chat

Embedding Models (use --embedding-model):
  • all-MiniLM-L6-v2 (default, fast, 384 dim)
  • all-mpnet-base-v2 (best quality, 768 dim)
  • paraphrase-multilingual-MiniLM-L12-v2 (multilingual)
        """
    )
    parser.add_argument("--index", action="store_true", 
                       help="Index all documents in the data folder")
    parser.add_argument("--force", action="store_true", 
                       help="Force reindex even if unchanged")
    parser.add_argument("--question", type=str, 
                       help="Ask a single question")
    parser.add_argument("--stats", action="store_true", 
                       help="Show database statistics")
    parser.add_argument("--embedding-model", type=str, 
                       default="all-MiniLM-L6-v2",
                       help="Sentence transformer model name")
    
    args = parser.parse_args()
    
    # Load environment variables
    load_dotenv()
    
    print("🚀 Initializing CyberSamatha RAG System...\n")
    
    # Initialize RAG system (Gemini is optional now)
    try:
        rag = CyberSamathaRAG(embedding_model=args.embedding_model)
        print("\n✅ System initialized successfully!\n")
    except Exception as e:
        print(f"❌ Failed to initialize RAG system: {e}")
        return
    
    # Show stats if requested
    if args.stats:
        stats = rag.get_collection_stats()
        print("="*60)
        print("📊 Vector Database Statistics")
        print("="*60)
        print(f"Total chunks:      {stats.get('total_chunks', 0)}")
        print(f"Collection:        {stats.get('collection_name', 'N/A')}")
        print(f"Embedding model:   {stats.get('embedding_model', 'N/A')}")
        print(f"AI answers:        {'✅ Enabled' if stats.get('gemini_enabled') else '❌ Disabled (set GEMINI_API_KEY)'}")
        print(f"Database path:     {stats.get('database_path', 'N/A')}")
        print("="*60)
        return
    
    # Index documents if requested
    if args.index or args.force:
        print("🚀 Starting document indexing with local embeddings...\n")
        rag.index_documents(force_reindex=args.force)
        # Don't exit after indexing - allow chat to start
    
    # Handle single question or interactive mode
    if args.question:
        answer = rag.ask_question(args.question)
        print(f"\n💬 **Question:** {args.question}")
        print(f"\n🤖 **Answer:**\n{answer}")
    else:
        # If we indexed and no question was provided, ask if user wants to chat
        if args.index or args.force:
            print("\n" + "="*50)
            response = input("🎯 Start interactive chat? (y/n): ").strip().lower()
            if response in ['y', 'yes', '']:
                rag.interactive_chat()
        else:
            rag.interactive_chat()

if __name__ == "__main__":
    main()