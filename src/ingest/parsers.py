import os
import json
import yaml
import hashlib
from typing import List

try:
    import PyPDF2
    from docx import Document
    import pptx
except ImportError as e:
    print(f"❌ Missing required parsing package: {e}")

class DocumentParser:
    @staticmethod
    def read_file(file_path: str) -> str:
        """Read various file formats and extract text"""
        try:
            if file_path.endswith('.pdf'):
                return DocumentParser._read_pdf(file_path)
            elif file_path.endswith('.docx'):
                return DocumentParser._read_docx(file_path)
            elif file_path.endswith('.pptx'):
                return DocumentParser._read_pptx(file_path)
            elif file_path.endswith('.json'):
                return DocumentParser._read_json(file_path)
            elif file_path.endswith(('.yaml', '.yml')):
                return DocumentParser._read_yaml(file_path)
            else:  # .txt, .md, etc.
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except Exception as e:
            print(f"❌ Error reading {file_path}: {str(e)}")
            return ""

    @staticmethod
    def _read_pdf(file_path: str) -> str:
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
        except Exception as e:
            print(f"❌ Error reading PDF {file_path}: {e}")
        return text

    @staticmethod
    def _read_docx(file_path: str) -> str:
        try:
            doc = Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            print(f"❌ Error reading DOCX {file_path}: {e}")
            return ""

    @staticmethod
    def _read_pptx(file_path: str) -> str:
        text = ""
        try:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"❌ Error reading PPTX {file_path}: {e}")
        return text

    @staticmethod
    def _read_json(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        except Exception as e:
            print(f"❌ Error reading JSON {file_path}: {e}")
            return ""

    @staticmethod
    def _read_yaml(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
                return yaml.dump(data, default_flow_style=False)
        except Exception as e:
            print(f"❌ Error reading YAML {file_path}: {e}")
            return ""

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        if not text or len(text.strip()) == 0:
            return []
            
        if len(text) <= chunk_size:
            return [text.strip()]
        
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            if end < len(text):
                last_period = text.rfind('. ', start, end)
                last_newline = text.rfind('\n', start, end)
                
                if last_period != -1 and last_period > start + chunk_size // 2:
                    end = last_period + 1
                elif last_newline != -1 and last_newline > start + chunk_size // 2:
                    end = last_newline + 1
                    
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start = end - chunk_overlap
            
            if start >= len(text):
                break
        
        return chunks

    @staticmethod
    def get_file_hash(file_path: str) -> str:
        try:
            hasher = hashlib.md5()
            with open(file_path, 'rb') as f:
                buf = f.read()
                hasher.update(buf)
            return hasher.hexdigest()
        except Exception as e:
            print(f"❌ Error hashing file {file_path}: {e}")
            return ""
